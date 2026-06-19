"""Generate Hawkeye+ executive briefing as a PPTX deck.

Run:
    python docs/build_pptx.py

Output:
    docs/Hawkeye_Plus_Overview.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).resolve().parent / "Hawkeye_Plus_Overview.pptx"

# 16:9 slide dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PRIMARY = RGBColor(0x0F, 0x76, 0x6E)        # teal
PRIMARY_DARK = RGBColor(0x0A, 0x4E, 0x49)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)         # amber
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=TEXT_DARK, align=None,
                font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_header(slide, title, subtitle=None, *, kicker=None):
    # Accent bar
    add_rect(slide, Inches(0.5), Inches(0.5), Inches(0.12), Inches(0.7), PRIMARY)
    if kicker:
        add_textbox(slide, Inches(0.75), Inches(0.45), Inches(10), Inches(0.3),
                    kicker.upper(), font_size=11, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(0.75), Inches(0.7), Inches(12), Inches(0.7),
                title, font_size=28, bold=True, color=PRIMARY_DARK)
    if subtitle:
        add_textbox(slide, Inches(0.75), Inches(1.25), Inches(12), Inches(0.45),
                    subtitle, font_size=14, color=TEXT_MUTED)


def add_footer(slide, page_no, total):
    add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), BG_LIGHT)
    add_textbox(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.3),
                "Hawkeye+ AutoTrain Platform — Executive Briefing",
                font_size=10, color=TEXT_MUTED)
    add_textbox(slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.3),
                f"{page_no} / {total}", font_size=10, color=TEXT_MUTED)


def add_bullets(slide, left, top, width, height, items, *,
                font_size=14, color=TEXT_DARK, bullet_color=PRIMARY,
                line_spacing=1.25):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        # item can be a string or (bold_lead, rest)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        bullet = p.add_run()
        bullet.text = "● "
        bullet.font.size = Pt(font_size)
        bullet.font.color.rgb = bullet_color
        bullet.font.bold = True
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, title, body_lines, *,
             title_color=PRIMARY_DARK, accent=PRIMARY):
    # Card background
    card = add_rect(slide, left, top, width, height, WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    # Top accent strip
    add_rect(slide, left, top, width, Inches(0.08), accent)
    # Title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15),
                width - Inches(0.4), Inches(0.4),
                title, font_size=15, bold=True, color=title_color)
    # Body
    bullet_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.6),
        width - Inches(0.4), height - Inches(0.7))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.line_spacing = 1.2
        dot = p.add_run()
        dot.text = "• "
        dot.font.size = Pt(11)
        dot.font.color.rgb = accent
        dot.font.bold = True
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.color.rgb = TEXT_DARK


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

TOTAL_SLIDES = 17  # cover + 12 main + divider + 4 backup = 18; we'll set dynamically


def slide_cover(prs):
    slide = add_blank_slide(prs)
    # Full background
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY_DARK)
    # Decorative band
    add_rect(slide, Inches(0), Inches(5.6), SLIDE_W, Inches(0.15), ACCENT)
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
                "EXECUTIVE BRIEFING", font_size=14, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.8), Inches(2.1), Inches(12), Inches(1.3),
                "Hawkeye+ AutoTrain Platform",
                font_size=48, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.4), Inches(12), Inches(0.8),
                "An End-to-End, Visualized Training Platform for Object Detection",
                font_size=20, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(0.5),
                "From Data to Deployable Model — In One Place",
                font_size=14, color=ACCENT)
    return slide


def slide_background(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Project Background",
               "Why We Built Hawkeye+", kicker="Slide 1")

    add_textbox(slide, Inches(0.75), Inches(1.85), Inches(12), Inches(0.4),
                "Industry Context", font_size=16, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(0.9), Inches(2.2), Inches(12), Inches(1.5), [
        "Custom object-detection models are increasingly demanded by industrial inspection and defect detection.",
        "Mainstream frameworks (e.g., PaddleDetection) are powerful but command-line driven, with hand-edited YAML and limited training visibility.",
        "Data annotation, format conversion, training, evaluation, and deployment are fragmented across disjointed tools.",
    ], font_size=13)

    add_textbox(slide, Inches(0.75), Inches(3.8), Inches(12), Inches(0.4),
                "Internal Pain Points", font_size=16, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(0.9), Inches(4.15), Inches(12), Inches(1.8), [
        "Algorithm engineers spend significant time on repetitive setup rather than model improvement.",
        "Non-algorithm staff (annotators, QA) cannot directly participate in the model iteration loop.",
        "Shared GPU servers lack governance, causing resource conflicts and unclear accountability.",
        "Training runs are opaque; experiments are hard to track and reproduce.",
    ], font_size=13)

    add_textbox(slide, Inches(0.75), Inches(6.05), Inches(12), Inches(0.4),
                "Our Goal", font_size=16, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.9), Inches(6.4), Inches(12), Inches(0.7),
                "Deliver a web-based, visualized, multi-user training platform that turns model development "
                "into a configure-and-click experience — with full transparency and control.",
                font_size=13, color=TEXT_DARK)


def slide_overview(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Platform Overview & Key Differentiators",
               "End-to-end, fully visualized training platform built on PaddleDetection",
               kicker="Slide 2")

    cards = [
        ("End-to-End Coverage",
         ["Projects, datasets, models, training, monitoring, validation in one place."]),
        ("Zero Command Line",
         ["GUI-driven operations; YAML configurations editable / importable visually."]),
        ("Real-Time Visibility",
         ["Live charts for loss, learning rate, mAP, GPU utilization, and ETA."]),
        ("Multi-User Collaboration",
         ["Role-based access (Admin / User), isolation, full activity audit log."]),
        ("Built-In Data Preparation",
         ["Integrated annotation tool plus one-click Labelme to COCO conversion."]),
        ("Deployment-Ready Output",
         ["TensorRT model export and download directly from the platform."]),
    ]
    # 3 columns x 2 rows
    card_w = Inches(4.0)
    card_h = Inches(1.55)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.7)
    start_y = Inches(2.0)
    for idx, (title, body) in enumerate(cards):
        r = idx // 3
        c = idx % 3
        left = start_x + c * (card_w + gap_x)
        top = start_y + r * (card_h + gap_y)
        add_card(slide, left, top, card_w, card_h, title, body)

    # Roles strip
    add_rect(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(1.3), BG_LIGHT,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    add_textbox(slide, Inches(0.9), Inches(5.6), Inches(12), Inches(0.4),
                "User Roles", font_size=14, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(0.9), Inches(5.95), Inches(12), Inches(0.85), [
        ("Administrator — ", "System settings, user management, data annotation, plus all standard features."),
        ("User — ", "Full self-service workflow: project → dataset → model → training → validation."),
    ], font_size=12)


def slide_workflow(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "End-to-End Workflow at a Glance",
               "Seven sequential, iterable steps from data to deployable model",
               kicker="Slide 3")

    steps = [
        ("1", "Projects"),
        ("2", "Datasets"),
        ("3", "Models"),
        ("4", "Configurations"),
        ("5", "Jobs"),
        ("6", "Monitoring"),
        ("7", "Validation"),
    ]
    n = len(steps)
    total_w = Inches(12.0)
    box_w = Inches(1.5)
    box_h = Inches(1.0)
    start_x = Inches(0.7)
    top = Inches(3.0)
    gap = (total_w - box_w * n) / (n - 1)
    for i, (num, label) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_rect(slide, x, top, box_w, box_h, PRIMARY)
        add_textbox(slide, x, top + Inches(0.1), box_w, Inches(0.4),
                    num, font_size=22, bold=True, color=WHITE, align=2)  # center
        add_textbox(slide, x, top + Inches(0.55), box_w, Inches(0.4),
                    label, font_size=12, bold=True, color=WHITE, align=2)
        # Arrow between
        if i < n - 1:
            ax = x + box_w + Inches(0.05)
            ay = top + Inches(0.4)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Inches(0.1), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    # Support modules row
    add_textbox(slide, Inches(0.7), Inches(4.5), Inches(12), Inches(0.4),
                "Supporting Modules", font_size=14, bold=True, color=PRIMARY)
    support = [
        ("Dashboard", "Global overview, GPU & job status in real time"),
        ("Data Annotation (Admin)", "Internal labeling and dataset preparation"),
        ("User Management (Admin)", "Roles, permissions, audit log"),
        ("Settings (Admin)", "System paths, GPU mappings, environment check"),
    ]
    card_w = Inches(2.95)
    card_h = Inches(1.4)
    gap_x = Inches(0.13)
    for i, (t, b) in enumerate(support):
        left = Inches(0.7) + i * (card_w + gap_x)
        add_card(slide, left, Inches(4.95), card_w, card_h, t, [b], accent=ACCENT)

    # Highlight strip
    add_rect(slide, Inches(0.7), Inches(6.55), Inches(12), Inches(0.55), BG_LIGHT,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    add_textbox(slide, Inches(0.9), Inches(6.62), Inches(12), Inches(0.45),
                "Persistent status indicator (System Ready / GPU Occupied / Training / Environment Error) is visible across the platform.",
                font_size=12, color=TEXT_DARK)


def two_column_step_slide(prs, *, slide_no, kicker, title, subtitle,
                          purpose, capabilities, value_items):
    slide = add_blank_slide(prs)
    add_header(slide, title, subtitle, kicker=kicker)

    # Purpose strip
    add_rect(slide, Inches(0.7), Inches(1.85), Inches(12), Inches(0.65),
             BG_LIGHT, line=RGBColor(0xE5, 0xE7, 0xEB))
    add_textbox(slide, Inches(0.9), Inches(1.92), Inches(2), Inches(0.5),
                "PURPOSE", font_size=11, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(2.5), Inches(1.92), Inches(10.2), Inches(0.5),
                purpose, font_size=13, color=TEXT_DARK)

    # Two columns
    col_w = Inches(5.9)
    col_h = Inches(4.2)
    left_x = Inches(0.7)
    right_x = Inches(6.75)
    top = Inches(2.75)

    # Capabilities
    add_rect(slide, left_x, top, col_w, col_h, WHITE,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(slide, left_x, top, col_w, Inches(0.5), PRIMARY)
    add_textbox(slide, left_x + Inches(0.2), top + Inches(0.07),
                col_w - Inches(0.4), Inches(0.4),
                "Capabilities", font_size=14, bold=True, color=WHITE)
    add_bullets(slide, left_x + Inches(0.2), top + Inches(0.65),
                col_w - Inches(0.4), col_h - Inches(0.75), capabilities,
                font_size=12)

    # Business Value
    add_rect(slide, right_x, top, col_w, col_h, WHITE,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(slide, right_x, top, col_w, Inches(0.5), ACCENT)
    add_textbox(slide, right_x + Inches(0.2), top + Inches(0.07),
                col_w - Inches(0.4), Inches(0.4),
                "Business Value", font_size=14, bold=True, color=WHITE)
    add_bullets(slide, right_x + Inches(0.2), top + Inches(0.65),
                col_w - Inches(0.4), col_h - Inches(0.75), value_items,
                font_size=12, bullet_color=ACCENT)


def slide_step1(prs):
    two_column_step_slide(
        prs,
        slide_no=4,
        kicker="Slide 4 — Step 1",
        title="Project Management",
        subtitle="The top-level container that organizes all training assets",
        purpose="Provide an organized, isolated container for every training initiative.",
        capabilities=[
            "Create, edit, search, and delete projects with unique names and descriptions.",
            "All datasets, models, configurations, and jobs are scoped to a project.",
            "Card-based listing for quick navigation across initiatives.",
        ],
        value_items=[
            ("Clear organization. ", "Every asset is traceable to a business initiative."),
            ("Strong isolation. ", "Products, versions, and customers do not interfere with each other."),
            ("Lifecycle control. ", "Legacy projects can be archived rather than lost."),
        ],
    )


def slide_step2(prs):
    two_column_step_slide(
        prs,
        slide_no=5,
        kicker="Slide 5 — Step 2",
        title="Dataset Management",
        subtitle="Ensure training data is available, healthy, and ready to use",
        purpose="Answer where training data comes from and verify its quality before training.",
        capabilities=[
            "Chunked web upload (large files, resume-on-failure).",
            "Server-side import of existing datasets with auto-detected paths.",
            "Supported formats: COCO and Labelme.",
            "One-click Labelme to COCO conversion with configurable train / val / test split.",
            "Dataset insight: annotation counts, class distribution, sample preview with bounding boxes.",
        ],
        value_items=[
            ("Large-file friendly. ", "Chunked transfer with resume support."),
            ("Format agnostic. ", "Built-in conversion eliminates external tooling."),
            ("Quality up-front. ", "Visual inspection catches issues before training starts."),
        ],
    )


def slide_step3(prs):
    two_column_step_slide(
        prs,
        slide_no=6,
        kicker="Slide 6 — Step 3",
        title="Model Configuration",
        subtitle="Decide which network architecture will learn the task",
        purpose="Allow users at any skill level to choose or define the model architecture.",
        capabilities=[
            "Default templates from PaddleDetection.",
            "Reusable user configurations saved by the team.",
            "Custom YAML for advanced users.",
            "Each model is named and bound to a project.",
        ],
        value_items=[
            ("Flexible by skill level. ", "Beginners use templates, experts customize."),
            ("Knowledge accumulation. ", "Proven configurations become team assets."),
            ("Full transparency. ", "Underlying YAML is always inspectable."),
        ],
    )


def slide_step4(prs):
    two_column_step_slide(
        prs,
        slide_no=7,
        kicker="Slide 7 — Step 4",
        title="Training Configuration",
        subtitle="Define how the model will be trained — hyperparameters and schedules",
        purpose="Capture and reuse training recipes independent of model and data.",
        capabilities=[
            "Parameter coverage: epochs, batch size, learning rate, scheduler, warmup, worker count, and more.",
            "Same three configuration sources: defaults, user history, custom YAML.",
            "Configurations are independent assets, reusable across jobs.",
        ],
        value_items=[
            ("Reusable recipes. ", "Proven setups power many experiments."),
            ("Decoupled design. ", "Any model can be paired with any configuration and dataset."),
        ],
    )


def slide_step5(prs):
    two_column_step_slide(
        prs,
        slide_no=8,
        kicker="Slide 8 — Step 5",
        title="Training Jobs",
        subtitle="Turn data + model + configuration into actual GPU work",
        purpose="Execute training under controlled, observable, and recoverable conditions.",
        capabilities=[
            "Job submission with selectable training configuration and explicit GPU assignment (single or multi-GPU).",
            "Training options: AMP (mixed precision) and VDL (VisualDL logging).",
            "Full lifecycle: Pending → Running → Completed / Failed / Stopped.",
            "Actions: start, stop, resume from checkpoint, delete, view live logs.",
        ],
        value_items=[
            ("Resource control. ", "Explicit GPU allocation prevents conflicts."),
            ("Resilient. ", "Long runs can be stopped and resumed from the latest checkpoint."),
            ("At-a-glance status. ", "Color-coded badges make queue health obvious."),
        ],
    )


def slide_step6(prs):
    two_column_step_slide(
        prs,
        slide_no=9,
        kicker="Slide 9 — Step 6",
        title="Training Monitoring",
        subtitle="Eliminate the black box during training",
        purpose="Provide real-time visibility, full traceability, and storage discipline during training.",
        capabilities=[
            "Live charts: loss, learning rate, mAP, iterations per second, ETA.",
            "Per-GPU resource monitoring (memory and utilization).",
            "Streaming log view with error highlighting and search.",
            "Checkpoint management: auto-save, best-model tracking, resume, cleanup.",
        ],
        value_items=[
            ("Real-time insight. ", "Issues are caught early, not after wasted compute."),
            ("Auditability. ", "Full historical logs and checkpoints are retained."),
            ("Storage efficiency. ", "Automated cleanup of obsolete checkpoints."),
        ],
    )


def slide_step7(prs):
    two_column_step_slide(
        prs,
        slide_no=10,
        kicker="Slide 10 — Step 7",
        title="Validation, Inference & Export",
        subtitle="Take a trained model from completed to evaluated, used, and deployable",
        purpose="Close the loop from training to deployment without leaving the platform.",
        capabilities=[
            "Select any completed training job and a specific checkpoint.",
            "Evaluation: mAP@0.5 / 0.75 / 0.5:0.95, AR@1 / 10 / 100, accuracy by object size (S/M/L), interactive charts.",
            "Inference: single image or batch folder, with annotated output images.",
            "TensorRT export: one-click export and download for deployment.",
            "Validation history: complete record of every evaluation and inference run.",
        ],
        value_items=[
            ("Closes the loop. ", "No separate deployment toolchain required."),
            ("Model comparison. ", "Different checkpoints can be evaluated side by side."),
            ("Visual evidence. ", "Charts and annotated images make quality intuitive for non-experts."),
        ],
    )


def slide_support(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Supporting Modules",
               "Cross-cutting features that complement the core workflow",
               kicker="Slide 11")

    cards = [
        ("Dashboard",
         ["Global overview of projects, jobs, GPUs, and recent activity."]),
        ("Annotation (Admin)",
         ["Built-in image annotation tool.",
          "Class selection, drag-to-draw, keyboard navigation.",
          "Auto-save every 30 seconds."]),
        ("User Management (Admin)",
         ["Create / edit / disable / delete users.",
          "Assign roles, reset passwords.",
          "Per-user activity log for audit."]),
        ("Settings (Admin)",
         ["Python and PaddleDetection paths.",
          "Per-GPU Python environment mapping.",
          "One-click environment health check."]),
    ]
    card_w = Inches(2.95)
    card_h = Inches(3.5)
    gap_x = Inches(0.13)
    for i, (t, body) in enumerate(cards):
        left = Inches(0.7) + i * (card_w + gap_x)
        add_card(slide, left, Inches(2.0), card_w, card_h, t, body)

    add_rect(slide, Inches(0.7), Inches(5.85), Inches(12), Inches(0.85),
             BG_LIGHT, line=RGBColor(0xE5, 0xE7, 0xEB))
    add_textbox(slide, Inches(0.9), Inches(5.95), Inches(12), Inches(0.65),
                "Health check verifies Python, CUDA, PaddleDetection, and disk space; results are surfaced as a live status badge in the global header.",
                font_size=12, color=TEXT_DARK)


def slide_summary(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Summary & Value Delivered",
               "A single web platform for the full object-detection training lifecycle",
               kicker="Slide 12")

    cards = [
        ("Value to Algorithm Engineers",
         ["Removes repetitive setup work.",
          "Focus shifts to data and modeling.",
          "Faster experiment cadence and iteration."]),
        ("Value to Non-Algorithm Staff",
         ["Lowers the barrier to participation.",
          "Annotators and QA contribute directly.",
          "Common language across functions."]),
        ("Value to Leadership & Operations",
         ["Visible GPU and job resource usage.",
          "Controlled, auditable processes.",
          "Traceable outputs from data to deployment."]),
    ]
    card_w = Inches(3.95)
    card_h = Inches(3.7)
    gap_x = Inches(0.2)
    for i, (t, body) in enumerate(cards):
        left = Inches(0.7) + i * (card_w + gap_x)
        add_card(slide, left, Inches(2.0), card_w, card_h, t, body, accent=PRIMARY)

    # Closing line
    add_rect(slide, Inches(0.7), Inches(6.05), Inches(12), Inches(0.85),
             PRIMARY_DARK)
    add_textbox(slide, Inches(0.9), Inches(6.15), Inches(12), Inches(0.65),
                "From data to deployable model — in one platform, with transparency, control, and collaboration.",
                font_size=14, bold=True, color=WHITE)


def slide_backup_divider(prs):
    slide = add_blank_slide(prs)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY_DARK)
    add_rect(slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.1), ACCENT)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(12), Inches(0.5),
                "BACKUP", font_size=16, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(12), Inches(1.2),
                "Technical Implementation",
                font_size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.7), Inches(12), Inches(0.6),
                "Architecture, data model, real-time training channel, and engineering details",
                font_size=16, color=WHITE)


def slide_backup_architecture(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Overall Technical Architecture",
               "Web application + dedicated training microservice",
               kicker="Backup 1")

    cards = [
        ("Frontend",
         ["Next.js 16 (App Router) + React 19 + TypeScript 5.",
          "Tailwind CSS 4 + shadcn/ui (Radix UI).",
          "Zustand + TanStack Query for state and data.",
          "React Hook Form + Zod for forms and validation.",
          "Recharts, Framer Motion, next-intl, next-themes."]),
        ("Backend",
         ["Next.js API Routes (40+ endpoints under src/app/api/).",
          "NextAuth.js with role-based middleware.",
          "Prisma ORM over SQLite."]),
        ("Training Microservice",
         ["Bun runtime with Socket.io on port 3003.",
          "Spawns training subprocesses.",
          "Parses PaddleDetection logs.",
          "Streams progress in real time."]),
        ("Deployment",
         ["next build produces a standalone bundle.",
          "Bun serves the application.",
          "Caddy reverse proxy.",
          "start.vbs for one-click startup on Windows."]),
    ]
    card_w = Inches(2.95)
    card_h = Inches(4.6)
    gap_x = Inches(0.13)
    for i, (t, body) in enumerate(cards):
        left = Inches(0.7) + i * (card_w + gap_x)
        add_card(slide, left, Inches(2.0), card_w, card_h, t, body)


def slide_backup_datamodel(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Data Model (Prisma Schema)",
               "Project-rooted ownership with audit-ready logging",
               kicker="Backup 2")

    add_textbox(slide, Inches(0.75), Inches(1.95), Inches(12), Inches(0.4),
                "Core Entities", font_size=15, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(0.9), Inches(2.3), Inches(12), Inches(2.5), [
        "User, SystemConfig.",
        "Project → Dataset / Model / TrainingConfig.",
        "TrainingJob → TrainingLog.",
        "ValidationJob.",
        "GpuMetric, ActivityLog.",
    ], font_size=13)

    add_textbox(slide, Inches(0.75), Inches(4.7), Inches(12), Inches(0.4),
                "Design Notes", font_size=15, bold=True, color=ACCENT)
    add_bullets(slide, Inches(0.9), Inches(5.05), Inches(12), Inches(2), [
        "All training assets are rooted at Project for clean ownership.",
        "TrainingLog persists per-epoch / per-iteration metrics for later analysis.",
        "ActivityLog records user actions to support auditing and compliance.",
    ], font_size=13, bullet_color=ACCENT)


def slide_backup_realtime(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Real-Time Training Channel (Socket.io)",
               "Bidirectional events drive live progress, logs, and lifecycle",
               kicker="Backup 3")

    # Two columns: events
    col_w = Inches(5.9)
    col_h = Inches(3.6)
    left_x = Inches(0.7)
    right_x = Inches(6.75)
    top = Inches(2.0)

    add_rect(slide, left_x, top, col_w, col_h, WHITE,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(slide, left_x, top, col_w, Inches(0.5), PRIMARY)
    add_textbox(slide, left_x + Inches(0.2), top + Inches(0.07),
                col_w - Inches(0.4), Inches(0.4),
                "Client → Server", font_size=14, bold=True, color=WHITE)
    add_bullets(slide, left_x + Inches(0.2), top + Inches(0.65),
                col_w - Inches(0.4), col_h - Inches(0.75), [
                    "training:start",
                    "training:stop",
                    "training:status",
                    "training:subscribe",
                    "training:unsubscribe",
                ], font_size=13)

    add_rect(slide, right_x, top, col_w, col_h, WHITE,
             line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(slide, right_x, top, col_w, Inches(0.5), ACCENT)
    add_textbox(slide, right_x + Inches(0.2), top + Inches(0.07),
                col_w - Inches(0.4), Inches(0.4),
                "Server → Client", font_size=14, bold=True, color=WHITE)
    add_bullets(slide, right_x + Inches(0.2), top + Inches(0.65),
                col_w - Inches(0.4), col_h - Inches(0.75), [
                    "training:started",
                    "training:log",
                    "training:progress",
                    "training:epoch",
                    "training:complete / stopped / error",
                ], font_size=13, bullet_color=ACCENT)

    add_textbox(slide, Inches(0.75), Inches(5.85), Inches(12), Inches(0.4),
                "Log Pipeline", font_size=15, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(0.9), Inches(6.2), Inches(12), Inches(0.9), [
        "A dedicated parser (log-parser.ts) extracts epoch, iteration, loss, learning rate, ETA, and memory metrics.",
        "Parsed records are persisted to the database and streamed to subscribed clients.",
    ], font_size=12)


def slide_backup_details(prs):
    slide = add_blank_slide(prs)
    add_header(slide, "Key Engineering Details",
               "Highlights of the implementation that enable the platform's promises",
               kicker="Backup 4")

    items = [
        ("Multi-GPU, multi-environment. ",
         "SystemConfig.gpuPythonMappings binds each GPU to its own Python environment."),
        ("Large-file upload. ",
         "Chunked transfer with resume-on-failure for datasets."),
        ("Three configuration sources. ",
         "Default templates, user history, and custom YAML — unified as YAML strings in the database."),
        ("Job isolation. ",
         "Each training job has its own outputDir, weightsPath, and vdlLogDir."),
        ("Environment health check. ",
         "/api/system/environment-check verifies Python, PaddleDetection, and GPU status; surfaced as a live status badge."),
    ]
    add_bullets(slide, Inches(0.9), Inches(2.1), Inches(12), Inches(4.8),
                items, font_size=14, line_spacing=1.4)


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,            # cover
        slide_background,       # 1
        slide_overview,         # 2
        slide_workflow,         # 3
        slide_step1,            # 4
        slide_step2,            # 5
        slide_step3,            # 6
        slide_step4,            # 7
        slide_step5,            # 8
        slide_step6,            # 9
        slide_step7,            # 10
        slide_support,          # 11
        slide_summary,          # 12
        slide_backup_divider,   # backup divider
        slide_backup_architecture,  # backup 1
        slide_backup_datamodel,     # backup 2
        slide_backup_realtime,      # backup 3
        slide_backup_details,       # backup 4
    ]
    total = len(builders)
    for i, build_fn in enumerate(builders):
        build_fn(prs)
        # Add footer/page number to all but the cover and backup divider
        slide = prs.slides[i]
        if build_fn not in (slide_cover, slide_backup_divider):
            add_footer(slide, i, total - 1)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
